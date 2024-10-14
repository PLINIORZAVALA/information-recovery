import os
import unicodedata
import pandas as pd
from nltk.tokenize import word_tokenize
from nltk.corpus import stopwords
from nltk.stem import SnowballStemmer
from collections import Counter
import docx
import PyPDF2
import pptx
from openpyxl import load_workbook
from nltk import download
import re 
import math
from collections import Counter
from sklearn.metrics.pairwise import cosine_similarity
import numpy as np
from sklearn.metrics import confusion_matrix, classification_report

# Descargar recursos de NLTK
download('punkt')
download('stopwords')

# Funciones para leer archivos
def read_word(filepath):
    doc = docx.Document(filepath)
    return "\n".join([para.text for para in doc.paragraphs])

def read_pdf(filepath):
    with open(filepath, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        text = ""
        for page in reader.pages:
            text += page.extract_text()
        return text

def read_powerpoint(filepath):
    ppt = pptx.Presentation(filepath)
    text = ""
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def read_excel(filepath):
    workbook = load_workbook(filepath)
    text = ""
    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        for row in sheet.iter_rows(values_only=True):
            text += " ".join([str(cell) for cell in row if cell]) + "\n"
    return text

def create_folder():
    output_folder = 'PreprocesamientoSteps'
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)
    return output_folder
# Leer archivo con los steams
def read_steams_from_excel(file_path):
    df = pd.read_excel(file_path)
    steams = df['Termino'].astype(str).tolist()  # Aseguramos que todo sea cadena
    return steams


# Función hash personalizada para evitar colisiones
def custom_hash(steam, table_size):
    return hash(steam) % table_size

# Función para indexar (hashing) los steams con manejo de colisiones usando sondeo lineal
def index_steams():
    # Leer los steams del archivo generado en el proceso anterior (4DiccDeSteams.xlsx)
    file_path = "PreprocesamientoSteps/4DiccDeSteams.xlsx"
    if not os.path.exists(file_path):
        print(f"El archivo {file_path} no existe.")
        return

    steams = read_steams_from_excel(file_path)

    # Definir el tamaño de la tabla hash (puede ajustarse según la cantidad de datos)
    table_size = len(steams)

    # Crear una tabla hash vacía (None significa que la posición está libre)
    steam_hash_table = [None] * table_size

    # Función para insertar steams usando sondeo lineal
    def insert_with_linear_probing(steam):
        # Calcular el hash inicial
        index = custom_hash(steam, table_size)
        
        # Sondeo lineal: buscar el siguiente espacio disponible
        original_index = index
        while steam_hash_table[index] is not None:
            index = (index + 1) % table_size  # Avanzar al siguiente índice (circular)
            if index == original_index:  # Si volvemos al punto de inicio, la tabla está llena
                raise Exception("Tabla hash llena, no se puede insertar más elementos.")
        
        # Insertar el steam en el índice disponible
        steam_hash_table[index] = steam
        return index

    # Insertar cada steam en la tabla hash usando sondeo lineal
    steam_to_hash_mapping = []
    for steam in steams:
        hash_index = insert_with_linear_probing(steam)
        steam_to_hash_mapping.append((steam, hash_index))

    # Guardar la tabla en un archivo Excel
    output_folder = create_folder()
    df = pd.DataFrame(steam_to_hash_mapping, columns=['Steam', 'Hash'])
    output_file = os.path.join(output_folder, '5ListDiccIndex.xlsx')
    df.to_excel(output_file, index=False)
    print(f"Tabla hash guardada en {output_file}.")
    
    return steam_hash_table  # Devolvemos la tabla hash para su uso posterior



# Función para procesar un solo documento (sin mayúsculas ni stemming aún)
def process_document_basic(content, stop_words_upper):
    tokens = word_tokenize(content)
    tokens = [word for word in tokens if word.isalnum() and not word.isnumeric()]
    tokens = [unicodedata.normalize('NFKD', word).encode('ascii', 'ignore').decode('utf-8') for word in tokens]
    
    # Eliminar stopwords (pero sin mayúsculas aún)
    tokens = [token for token in tokens if token.upper() not in stop_words_upper]
    
    return tokens

def count_documents(corpus_dir):
    document_count = 0

    # Recorrer todos los archivos en el directorio
    for filename in os.listdir(corpus_dir):
        file_path = os.path.join(corpus_dir, filename)
        
        # Verificar si el archivo tiene una extensión que corresponde a documentos soportados
        if filename.endswith(('.txt', '.pdf', '.docx', '.xlsx', '.pptx')):
            document_count += 1

    return document_count

def tokenize_and_process_documents(corpus_dir):
    initial_dictionary = {}
    stemmer = SnowballStemmer("spanish")
    stop_words = set(word.upper() for word in stopwords.words('spanish'))  # Stopwords en mayúsculas
    output_folder = create_folder()  # Crear la carpeta para los archivos de salida

    document_word_counts = {}

    # Proceso 1: Crear el diccionario inicial
    for filename in os.listdir(corpus_dir):
        file_path = os.path.join(corpus_dir, filename)
        content = ""

        if filename.endswith(".txt"):
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
        elif filename.endswith(".pdf"):
            content = read_pdf(file_path)
        elif filename.endswith(".docx"):
            content = read_word(file_path)
        elif filename.endswith(".xlsx"):
            content = read_excel(file_path)
        elif filename.endswith(".pptx"):
            content = read_powerpoint(file_path)
        else:
            print(f"Ignorando archivo no soportado: {filename}")
            continue

        # Aplicar el preprocesamiento básico (sin mayúsculas ni stemming)
        processed_tokens = process_document_basic(content, stop_words)
        tokens = word_tokenize(content)
        tokens = [word for word in tokens if word.isalnum() and not word.isnumeric()]
        tokens = [unicodedata.normalize('NFKD', word).encode('ascii', 'ignore').decode('utf-8') for word in tokens]

        # Contar la frecuencia de palabras preprocesadas en cada documento
        doc_word_counts = {}
        for token in tokens:
            token_stemmed = stemmer.stem(token).upper()  # Aplicar stemming y convertir a mayúsculas
            if token_stemmed not in doc_word_counts:
                doc_word_counts[token_stemmed] = 1
            else:
                doc_word_counts[token_stemmed] += 1

        document_word_counts[filename] = doc_word_counts

        # Actualizar el diccionario inicial con el conteo de todas las palabras
        for token in processed_tokens:
            if token not in initial_dictionary:
                initial_dictionary[token] = 1
            else:
                initial_dictionary[token] += 1

    # Guardar el diccionario inicial en "1Diccionario.xlsx" (sin convertir a mayúsculas ni aplicar stemming)
    df_initial = pd.DataFrame(list(initial_dictionary.items()), columns=['Termino', 'Frecuencia'])
    df_initial.to_excel(os.path.join(output_folder, '1Diccionario.xlsx'), index=False)
    print('Diccionario inicial guardado en "1Diccionario.xlsx".')

    # Proceso 2: Convertir el diccionario inicial a mayúsculas
    upper_dictionary = {}
    for token in initial_dictionary:
        token_upper = token.upper()
        if token_upper not in upper_dictionary:
            upper_dictionary[token_upper] = initial_dictionary[token]
        else:
            upper_dictionary[token_upper] += initial_dictionary[token]

    # Guardar el diccionario con mayúsculas en "2DiccMayus.xlsx"
    df_upper = pd.DataFrame(list(upper_dictionary.items()), columns=['Termino', 'Frecuencia'])
    df_upper.to_excel(os.path.join(output_folder, '2DiccMayus.xlsx'), index=False)
    print('Diccionario con mayúsculas guardado en "2DiccMayus.xlsx".')

    # Proceso 3: Eliminar stopwords del diccionario con mayúsculas
    stop_words_upper = {word.upper() for word in stop_words}
    filtered_dictionary = {word: count for word, count in upper_dictionary.items() if word not in stop_words_upper}

    # Guardar el diccionario sin stopwords en "3DiccMSinStopWords.xlsx"
    df_filtered = pd.DataFrame(list(filtered_dictionary.items()), columns=['Termino', 'Frecuencia'])
    df_filtered.to_excel(os.path.join(output_folder, '3DiccMSinStopWords.xlsx'), index=False)
    print('Diccionario sin stop words guardado en "3DiccMSinStopWords.xlsx".')

    # Proceso 4: Aplicar stemming al diccionario
    stemmed_dictionary = {}
    for word, count in filtered_dictionary.items():
        stemmed_word = stemmer.stem(word).upper()
        if stemmed_word not in stemmed_dictionary:
            stemmed_dictionary[stemmed_word] = count
        else:
            stemmed_dictionary[stemmed_word] += count

    # Guardar el diccionario con stemming en "4DiccDeSteams.xlsx"
    df_stemmed = pd.DataFrame(list(stemmed_dictionary.items()), columns=['Termino', 'Frecuencia'])
    df_stemmed.to_excel(os.path.join(output_folder, '4DiccDeSteams.xlsx'), index=False)
    print('Diccionario con stemming guardado en "4DiccDeSteams.xlsx".')

    # Proceso 5: Tbla de hash
    index_steams()
    
    # Llamar a la función con el argumento correcto
    save_document_word_counts(document_word_counts, output_folder, stemmed_dictionary)
    print("Frecuencia de términos por documento guardada correctamente.")

    total_documents = count_documents(corpus_dir)
    calculate_tf(Frecuencia_Terminos, total_documents, output_folder)

    # Leer el archivo de conteo de documentos por término
    document_term_count_file = os.path.join(output_folder, 'Conteo_Documentos_Termino.xlsx')
    # Calcular el IDF
    calculate_idf(total_documents, document_term_count_file, output_folder, document_word_counts)

    # Al final del procesamiento, después de calcular TF e IDF
    calculate_tf_idf(Frecuencia_Terminos, os.path.join(output_folder, 'IDF.xlsx'), output_folder)
    
    # 

def save_document_word_counts(document_word_counts, output_folder, stemmed_dictionary):
    # Obtener todas las palabras únicas del diccionario con stemming
    all_words = set(stemmed_dictionary.keys())

    # Filtrar palabras que contienen números
    all_words = {word for word in all_words if not any(char.isdigit() for char in word)}

    all_words = sorted(list(all_words))  # Ordenar todas las palabras alfabéticamente
    
    # Crear un DataFrame con los nombres de los documentos en las filas y los términos en las columnas
    df = pd.DataFrame(index=list(document_word_counts.keys()), columns=all_words)

    # Llenar el DataFrame con las frecuencias
    for doc, counts in document_word_counts.items():
        for word in all_words:
            df.at[doc, word] = counts.get(word, 0)  # Poner el conteo o 0 si no aparece la palabra

    # Convertir todos los términos a mayúsculas en el DataFrame
    df.columns = [word.upper() for word in df.columns]

    # Guardar el DataFrame en un archivo Excel
    output_file = os.path.join(output_folder, 'Frecuencia_Terminos_Por_Documento.xlsx')
    df.to_excel(output_file, index=True)
    print(f"Frecuencia de términos por documento guardada en: {output_file}")

    # Crear un diccionario para contar en cuántos documentos aparece cada término
    term_document_count = {}

    for word in all_words:
        term_count_in_docs = (df[word] > 0).sum()  # Contar cuántos documentos tienen al menos una aparición del término
        term_document_count[word] = term_count_in_docs

    # Guardar el conteo de documentos por término en un archivo Excel
    df_doc_count = pd.DataFrame(list(term_document_count.items()), columns=['Termino', 'Documentos_Contienen_Termino'])
    output_doc_count_file = os.path.join(output_folder, 'Conteo_Documentos_Termino.xlsx')
    df_doc_count.to_excel(output_doc_count_file, index=False)
    print(f"Conteo de documentos por término guardado en: {output_doc_count_file}")

Frecuencia_Terminos = 'PreprocesamientoSteps/Frecuencia_Terminos_Por_Documento.xlsx'  # Asegúrate de dar la ruta correcta

def calculate_tf(Frecuencia_Terminos, total_documents, output_folder):
    # Leer el archivo Excel existente
    df = pd.read_excel(Frecuencia_Terminos, index_col=0)  # Usa la primera columna como índice (nombres de documentos)
    
    # Dividir cada valor por el total de documentos para obtener TF
    df_tf = df.div(total_documents)

    # Guardar el DataFrame en un nuevo archivo Excel
    output_file = os.path.join(output_folder, 'TF.xlsx')
    df_tf.to_excel(output_file, index=True)
    print(f"TF guardado en: {output_file}")

def calculate_idf(total_documents, document_term_count_file, output_folder, document_word_counts):
    # Leer el archivo de conteo de documentos por término
    df_term_doc_count = pd.read_excel(document_term_count_file)

    # Calcular el IDF (basado en el total de documentos)
    df_term_doc_count['IDF'] = df_term_doc_count['Documentos_Contienen_Termino'].apply(
        lambda doc_count: math.log2(total_documents / doc_count) if doc_count > 0 else 0
    )
    
    # Crear un DataFrame con los nombres de los documentos en las filas y los términos en las columnas
    df_idf = pd.DataFrame(index=list(document_word_counts.keys()), columns=df_term_doc_count['Termino'])

    # Llenar cada columna con su valor de IDF correspondiente
    for term in df_idf.columns:
        if term in df_term_doc_count['Termino'].values:
            idf_value = df_term_doc_count.loc[df_term_doc_count['Termino'] == term, 'IDF'].values[0]
            df_idf[term] = idf_value  # Rellenar toda la columna con el mismo valor de IDF
        else:
            df_idf[term] = 0  # Si no existe el término, llenarlo con 0

    # Guardar el DataFrame en un archivo Excel
    output_file = os.path.join(output_folder, 'IDF.xlsx')
    df_idf.to_excel(output_file, index=True)  # Guardamos el DataFrame con los nombres de los documentos como índice
    
    print(f"IDF guardado en: {output_file}")

def calculate_tf_idf(tf_file, idf_file, output_folder):
    # Leer los archivos de TF y IDF
    df_tf = pd.read_excel(tf_file, index_col=0)  # Cargar TF y usar el nombre de documentos como índice
    df_idf = pd.read_excel(idf_file, header=0)  # Cargar IDF y usar la primera fila como encabezado

    # Obtener los valores de IDF (la única fila de valores)
    idf_values = df_idf.iloc[0].values  # Obtener la fila de IDF
    idf_terms = df_idf.columns  # Obtener los términos de IDF

    # Asegurarse de que los términos de TF y IDF coincidan
    common_terms = df_tf.columns.intersection(idf_terms)  # Términos comunes entre TF e IDF

    if len(common_terms) == 0:
        raise ValueError("No hay términos en común entre TF e IDF.")

    # Filtrar DF TF para mantener solo los términos comunes
    df_tf_common = df_tf[common_terms]

    # Filtrar los valores IDF para que coincidan con los términos comunes
    idf_values_common = [idf_values[df_idf.columns.get_loc(term)] for term in common_terms]

    # Multiplicar cada valor en df_tf_common por su correspondiente en idf_values_common
    df_tf_idf = df_tf_common.multiply(idf_values_common, axis=1)

    # Guardar el DataFrame de TF-IDF en un archivo Excel
    output_file = os.path.join(output_folder, '8MatrizTF-IDF.xlsx')
    df_tf_idf.to_excel(output_file, index=True)
    print(f"TF-IDF guardado en: {output_file}")

def create_query_vector(query, stemmed_dictionary):
    # Preprocesar la consulta
    tokens = word_tokenize(query)
    stemmer = SnowballStemmer("spanish")
    stop_words = set(word.upper() for word in stopwords.words('spanish'))
    
    # Normalizar acentos y quitar signos de puntuación
    processed_tokens = [
        unicodedata.normalize('NFKD', token).encode('ascii', 'ignore').decode('utf-8').upper()
        for token in tokens if token.isalpha() and token.upper() not in stop_words
    ]

    # Realizar stemming y convertir a mayúsculas
    processed_tokens = [stemmer.stem(token).upper() for token in processed_tokens]
    
    print("Tokens procesados después del filtrado y normalización:", processed_tokens)

    # Contar la frecuencia de los stems en la consulta
    token_counts = Counter(processed_tokens)

    # Crear el vector de consulta
    query_vector = np.zeros(len(stemmed_dictionary))

    for i, stem in enumerate(stemmed_dictionary):
        if stem in token_counts:
            query_vector[i] = token_counts[stem]  # Asignar la frecuencia del término

    return query_vector

def calculate_similarity(query_vector, tf_idf_matrix):
    
    query_vector_reshaped = np.array(query_vector).reshape(1, -1)
    similarities = cosine_similarity(query_vector_reshaped, tf_idf_matrix)
    return similarities.flatten()  # Retornar como vector 1D

######################## MATRIZ DE CONFUSION #############################
def clean_document_name(name):
    """Función para limpiar los nombres de los documentos eliminando caracteres especiales y extensiones."""
    # Eliminar caracteres especiales como guiones y guiones bajos, y eliminar la extensión .pdf
    name = re.sub(r'[-_]', ' ', name)  # Reemplazar guiones y guiones bajos por espacios
    name = re.sub(r'\.(pdf|txt|docx|pptx)$', '', name, flags=re.IGNORECASE)  # Eliminar extensión .pdf si está presente
    name = name.strip().lower()  # Quitar espacios en blanco y pasar a minúsculas
    return name

def etiquetas_reales(file_path, documentos):
    try:
        # Leer el archivo Excel
        df = pd.read_excel(file_path)

        # Verificar que las columnas necesarias estén presentes
        if 'Documento' not in df.columns or 'Relevancia' not in df.columns:
            raise ValueError("El archivo debe contener las columnas 'Documento' y 'Relevancia'.")

        # Quitar espacios en blanco en ambas columnas, normalizar a minúsculas y limpiar nombres de documentos
        df['Documento'] = df['Documento'].apply(clean_document_name)
        df['Relevancia'] = df['Relevancia'].str.strip()

        # Filtrar solo las etiquetas válidas ('R' y 'NR')
        df = df[df['Relevancia'].isin(['R', 'NR'])]

        # Imprimir DataFrame para depuración
        print("\nContenido del DataFrame con etiquetas reales (después de limpiar):")
        print(df)

        # Generar lista de etiquetas reales en el mismo orden que los documentos
        etiquetas_reales = []
        for doc in documentos:
            doc_normalizado = clean_document_name(doc)  # Normalizar documento con la misma función
            etiqueta = df.loc[df['Documento'] == doc_normalizado, 'Relevancia']

            # Si no se encuentra, mostrar advertencia para depuración
            if not etiqueta.empty:
                etiquetas_reales.append(etiqueta.values[0])
            else:
                print(f"Advertencia: No se encontró etiqueta para el documento '{doc_normalizado}'")
                etiquetas_reales.append(None)  # Si no se encuentra etiqueta

        # Imprimir las etiquetas reales generadas para depuración
        #print("\nEtiquetas reales en el orden de los documentos:")
        #print(etiquetas_reales)

        return etiquetas_reales

    except Exception as e:
        print(f"Error al leer el archivo {file_path}: {e}")
        return []

# Función para generar predicciones en base a las similitudes
def get_predicciones(similarities, documentos, threshold=0.005):
    predicciones = []

    # Iterar sobre las similitudes y los documentos
    for i, sim in enumerate(similarities):
        doc = documentos[i]  # Obtener el documento correspondiente
        
        # Si la similitud es mayor al umbral, lo clasificamos como "R" (Relevante)
        if sim > threshold:
            predicciones.append('R')
            print(f"Documento: {doc} - Similitud: {sim:.4f} - Predicción: R")
        else:
            predicciones.append('NR')
            print(f"Documento: {doc} - Similitud: {sim:.4f} - Predicción: NR")

    return predicciones

# Función para evaluar el modelo comparando las etiquetas reales con las predicciones
def evaluar_modelo(etiquetas_reales, predicciones, documentos):
    # Inicializar contadores para las métricas
    tp = 0  # Verdaderos Positivos
    tn = 0  # Verdaderos Negativos
    fp = 0  # Falsos Positivos
    fn = 0  # Falsos Negativos

    # Imprimir la comparación entre etiquetas reales y predicciones
    print("\nComparación entre documentos, etiquetas reales y predicciones:")
    
    for i, doc in enumerate(documentos):
        etiqueta_real = etiquetas_reales[i]
        prediccion = predicciones[i]

        print(f"Documento: '{doc}', Etiqueta real: '{etiqueta_real}', Predicción: '{prediccion}'")

        # Comparar predicción con la etiqueta real y actualizar contadores
        if etiqueta_real is not None:  # Evaluar solo si hay una etiqueta real
            if prediccion == 'R' and etiqueta_real == 'R':
                tp += 1
            elif prediccion == 'NR' and etiqueta_real == 'NR':
                tn += 1
            elif prediccion == 'R' and etiqueta_real == 'NR':
                fp += 1
            elif prediccion == 'NR' and etiqueta_real == 'R':
                fn += 1

    print(f"\nVerdaderos Positivos (TP): {tp}")
    print(f"Verdaderos Negativos (TN): {tn}")
    print(f"Falsos Positivos (FP): {fp}")
    print(f"Falsos Negativos (FN): {fn}")

    # Calcular las métricas de evaluación
    exactitud = (tp + tn) / (tp + tn + fp + fn) if (tp + tn + fp + fn) > 0 else 0
    precision = tp / (tp + fp) if (tp + fp) > 0 else 0
    sensibilidad = tp / (tp + fn) if (tp + fn) > 0 else 0
    tasa_error = (fp + fn) / (tp + tn + fp + fn) if (tp + tn + fp + fn) > 0 else 0
    f1_score = 2 * (precision * sensibilidad) / (precision + sensibilidad) if (precision + sensibilidad) > 0 else 0

    # Guardar las métricas en un archivo de texto
    with open('PreprocesamientoSteps/9-5Evaluacionmodelo.txt', 'w') as f:
        f.write(f"Exactitud: {exactitud:.2f}\n")
        f.write(f"Precisión: {precision:.2f}\n")
        f.write(f"Sensibilidad: {sensibilidad:.2f}\n")
        f.write(f"Tasa de error: {tasa_error:.2f}\n")
        f.write(f"F1-score: {f1_score:.2f}\n")

    return tp, tn, fp, fn

def generate_matriz_confusion(real_labels, predicted_labels):
    if isinstance(real_labels, dict):
        real_labels = list(real_labels.values())

    predicted_labels = list(predicted_labels)
    cm = confusion_matrix(real_labels, predicted_labels)
    print("Matriz de confusión:")
    print(cm)

    # Guardar la matriz de confusión en un archivo Excel
    cm_df = pd.DataFrame(cm, index=['NR', 'R'], columns=['NR', 'R'])
    cm_df.to_excel('PreprocesamientoSteps/9Evaluacionmodelo.xlsx', index=True)

    report = classification_report(real_labels, predicted_labels)
    print("Informe de clasificación:\n", report)

def preprocess_documents():
    corpus_dir = 'corpus'
    
    total_documents = count_documents(corpus_dir)
    print(f"Número total de documentos en el corpus: {total_documents}")

    if os.path.exists(corpus_dir):
        tokenize_and_process_documents(corpus_dir)
        
        tf_idf_df = pd.read_excel('PreprocesamientoSteps/8MatrizTF-IDF.xlsx', index_col=0)
        tf_idf_matrix = tf_idf_df.values
        
        # Obtener el vocabulario de stems
        stemmed_dictionary = list(tf_idf_df.columns)

        query = input("Introduce tu consulta: ")
        query_vector = create_query_vector(query, stemmed_dictionary)

        similarities = calculate_similarity(query_vector, tf_idf_matrix)
        
        results_df = pd.DataFrame({
            'Documento': tf_idf_df.index,
            'Similitud': similarities
        })
        
        results_df = results_df.sort_values(by='Similitud', ascending=False)
        print("Resultados ordenados por similitud:\n", results_df)

        # Documentos que serán evaluados
        documentos = results_df['Documento'].tolist()
        similitudes = results_df['Similitud'].tolist()

        #print("Documentos para evaluación:", documentos)  # Imprime los documentos que se evaluarán
        #print("Similitudes:", similitudes)

        # Obtener etiquetas reales en el orden de los documentos
        real_labels = etiquetas_reales('RN-NR.xlsx', documentos)

        # Obtener predicciones a partir de las similitudes
        predicciones = get_predicciones(similitudes, documentos)  # Ajuste aquí para pasar solo similitudes

        # Evaluar el modelo
        evaluar_modelo(real_labels, predicciones, documentos)

    else:
        print("El directorio 'Corpus' no existe.")

# Ejecutar el procesamiento
preprocess_documents()
