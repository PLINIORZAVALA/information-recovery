import tkinter as tk
from tkinter import ttk
import os
import pandas as pd
import nltk
import unicodedata
from nltk.tokenize import word_tokenize
from nltk.corpus import stopwords
from nltk.stem import PorterStemmer
import PyPDF2
import docx
from pptx import Presentation
from nltk.stem.snowball import SnowballStemmer

nltk.download('punkt')
nltk.download('stopwords')

def opendocument(link):
    try:
        if os.name == 'nt':  # Para Windows
            os.startfile(link)
    except Exception as e:
        print(f"Error al intentar abrir el documento: {e}")

def create_folder():
    folder_name = "PreprocesamientoSteps"
    if not os.path.exists(folder_name):
        os.makedirs(folder_name)
        print(f'Carpeta "{folder_name}" creada con éxito.')
    else:
        print(f'La carpeta "{folder_name}" ya existe.')
    return folder_name  # Devuelve el nombre de la carpeta creada

def read_pdf(file_path):
    content = ""
    with open(file_path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        for page in reader.pages:
            content += page.extract_text()
    return content

def read_word(file_path):
    doc = docx.Document(file_path)
    content = "\n".join([para.text for para in doc.paragraphs])
    return content

def read_excel(file_path):
    content = ""
    df = pd.read_excel(file_path)
    for column in df.columns:
        content += " ".join(df[column].astype(str)) + " "
    return content

def read_powerpoint(file_path):
    prs = Presentation(file_path)
    content = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content += shape.text + "\n"
    return content

def tokenize_and_process_documents(corpus_dir):
    initial_dictionary = {}

    stemmer = PorterStemmer()
    stop_words = set(word.upper() for word in stopwords.words('spanish'))  # Stopwords en mayúsculas

    output_folder = create_folder()  # Crear la carpeta para los archivos de salida

    for filename in os.listdir(corpus_dir):
        file_path = os.path.join(corpus_dir, filename)
        content = ""

        if filename.endswith(".txt"):
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
        elif filename.endswith(".pdf"):
            content = read_pdf(file_path)
        elif filename.endswith(".docx"):
            content = read_word(file_path)
        elif filename.endswith(".xlsx"):
            content = read_excel(file_path)
        elif filename.endswith(".pptx"):
            content = read_powerpoint(file_path)
        else:
            print(f"Ignorando archivo no soportado: {filename}")
            continue
    # Inicio del primer proceso
        tokens = word_tokenize(content)
        tokens = [word for word in tokens if word.isalnum() and not word.isnumeric()]
        tokens = [unicodedata.normalize('NFKD', word).encode('ascii', 'ignore').decode('utf-8') for word in tokens]

        for token in tokens:
            if token not in initial_dictionary:
                initial_dictionary[token] = 1
            else:
                initial_dictionary[token] += 1

    # Guardar el diccionario inicial en "1Diccionario.xlsx"
    df_initial = pd.DataFrame(list(initial_dictionary.items()), columns=['Termino', 'Frecuencia'])
    df_initial.to_excel(os.path.join(output_folder, '1Diccionario.xlsx'), index=False)
    print('Diccionario inicial guardado en "1Diccionario.xlsx".')


    # inicio del segundo proceso
    upper_dictionary = {}

    # Proceso de cambio a mayúsculas tomando como base el diccionario inicial
    for token in initial_dictionary:
     token_upper = token.upper()
     if token_upper not in upper_dictionary:
         upper_dictionary[token_upper] = initial_dictionary[token]
    else:
         upper_dictionary[token_upper] += initial_dictionary[token]  # Si ya existe, sumar la frecuencia

    # Guardar el diccionario con mayúsculas en "2DiccMayus.xlsx"
    df_upper = pd.DataFrame(list(upper_dictionary.items()), columns=['Termino', 'Frecuencia'])
    df_upper.to_excel(os.path.join(output_folder, '2DiccMayus.xlsx'), index=False)
    print('Diccionario con mayúsculas guardado en "2DiccMayus.xlsx".')

    #Proceso numero 3
    stop_words = set(stopwords.words('spanish'))
    stop_words_upper = {word.upper() for word in stop_words}
    filtered_dictionary = {word: count for word, count in upper_dictionary.items() if word not in stop_words_upper}

    # Guardar el diccionario sin stop words en "3DiccMSinStopWords.xlsx"
    df_filtered = pd.DataFrame(list(filtered_dictionary.items()), columns=['Termino', 'Frecuencia'])
    df_filtered.to_excel(os.path.join(output_folder, '3DiccMSinStopWords.xlsx'), index=False)
    print('Diccionario sin stop words guardado en "3DiccMSinStopWords.xlsx".')
    
    #Proceso numnero 4
    stemmer = SnowballStemmer("spanish")
    stemmed_dictionary = {}
    
    for word, count in filtered_dictionary.items():
        stemmed_word = stemmer.stem(word).upper()
        if stemmed_word not in stemmed_dictionary:
            stemmed_dictionary[stemmed_word] = count
        else:
            stemmed_dictionary[stemmed_word] += count

    # Guardar el diccionario con stemming en "4DiccDeSteams.xlsx"
    df_stemmed = pd.DataFrame(list(stemmed_dictionary.items()), columns=['Termino', 'Frecuencia'])
    df_stemmed.to_excel(os.path.join(output_folder, '4DiccDeSteams.xlsx'), index=False)
    print('Diccionario con stemming guardado en "4DiccDeSteams.xlsx".')

def preprocess_documents():
    corpus_dir = 'Corpus'
    if os.path.exists(corpus_dir):
        tokenize_and_process_documents(corpus_dir)
    else:
        print("El directorio 'Corpus' no existe.")

def createapp():
    root = tk.Tk()
    root.title("Cronus")
    root.configure(bg='black')
    root.attributes('-fullscreen', True)

    style = ttk.Style()
    style.configure("TButton", font=("arial", 16))
    style.configure("TEntry", font=("arial", 16))

    search_frame = tk.Frame(root, bg='black')
    search_frame.pack(pady=20)

    search_entry = ttk.Entry(search_frame, width=30, style="TEntry")
    search_entry.pack(side=tk.LEFT, padx=5)

    search_button = ttk.Button(search_frame, text="🔍", style="TButton")
    search_button.pack(side=tk.LEFT)

    preprocess_button = ttk.Button(root, text="Preprocesado", command=preprocess_documents, style="TButton")
    preprocess_button.pack(pady=10)

    canvas = tk.Canvas(root, bg='gray')
    canvas.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

    scrollbar = ttk.Scrollbar(root, orient="vertical", command=canvas.yview)
    scrollbar.pack(side=tk.RIGHT, fill='y')

    canvas.configure(yscrollcommand=scrollbar.set)
    canvas.bind('<Configure>', lambda e: canvas.configure(scrollregion=canvas.bbox("all")))

    frame = tk.Frame(canvas, bg='gray')
    canvas.create_window((0, 0), window=frame, anchor='nw')

    corpus_dir = 'Corpus'
    if not os.path.exists(corpus_dir):
        os.makedirs(corpus_dir)

    documents = [f for f in os.listdir(corpus_dir) if os.path.isfile(os.path.join(corpus_dir, f))]

    for doc in documents:
        doc_path = os.path.join(corpus_dir, doc)
        link = tk.Label(frame, text=doc, fg="white", cursor="hand2", bg='gray', font=("arial", 16))
        link.pack(anchor='w', pady=2)
        link.bind("<Button-1>", lambda e, url=doc_path: opendocument(url))

    root.bind("<Escape>", lambda e: root.attributes("-fullscreen", False))
    root.mainloop()

if __name__ == "__main__":
    createapp()
