import os
import PyPDF2
import docx
import pandas as pd
from pptx import Presentation
import unicodedata
from nltk.tokenize import word_tokenize

# Funciones para leer diferentes tipos de archivos
def read_pdf(file_path):
    content = ""
    with open(file_path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        for page in reader.pages:
            content += page.extract_text()
    return content

def read_word(file_path):
    doc = docx.Document(file_path)
    content = "\n".join([para.text for para in doc.paragraphs])
    return content  

def read_excel(file_path):
    content = ""
    df = pd.read_excel(file_path)
    for column in df.columns:
        content += " ".join(df[column].astype(str)) + " "
    return content

def read_powerpoint(file_path):
    prs = Presentation(file_path)
    content = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content += shape.text + "\n"
    return content

# Función para tokenizar y procesar un documento
def tokenize_document(content):
    tokens = word_tokenize(content.upper())  # Tokenización en mayúsculas
    tokens = [word for word in tokens if word.isalpha() and len(word) > 1]  # Filtrar tokens alfabéticos y con más de 1 letra
    tokens = [unicodedata.normalize('NFKD', word).encode('ascii', 'ignore').decode('utf-8') for word in tokens]  # Normalización
    return tokens

# Función para listar steams desde un archivo Excel
def list_steams_from_excel(steams_dir):
    if not os.path.exists(steams_dir):
        print(f"El archivo {steams_dir} no existe.")
        return []
    
    df = pd.read_excel(steams_dir)
    
    if 'Steams' not in df.columns:
        print(f"La columna 'Steams' no existe en el archivo {steams_dir}.")
        return []
    
    steams = df['Steams'].astype(str).tolist()  # Convertimos todos los términos a cadenas
    return steams

# Procesar documentos, tokenizarlos y contar la cantidad de palabras tokenizadas
def process_and_tokenize_documents(corpus_dir, steams_dir, output_dir):
    steams = list_steams_from_excel(steams_dir)  # Cargar la lista de steams
    results = {}  # Diccionario para almacenar resultados por documento

    for idx, filename in enumerate(os.listdir(corpus_dir), start=1):
        file_path = os.path.join(corpus_dir, filename)
        content = ""
        
        if filename.endswith(".txt"):
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
        elif filename.endswith(".pdf"):
            content = read_pdf(file_path)
        elif filename.endswith(".docx"):
            content = read_word(file_path)
        elif filename.endswith(".xlsx"):
            content = read_excel(file_path)
        elif filename.endswith(".pptx"):
            content = read_powerpoint(file_path)
        else:
            print(f"Ignorando archivo no soportado: {filename}")
            continue

        # Tokenizar el contenido
        tokens = tokenize_document(content)

        # Filtrar steams encontrados en el documento y eliminar duplicados
        found_steams = sorted(set([token for token in tokens if token in steams]))

        # Almacenar los steams encontrados en el diccionario de resultados, incluso si está vacío
        results[f"D{idx}"] = found_steams if found_steams else [None]  # Colocar None si el documento no tiene steams
        print(f"Procesado archivo {filename} con {len(found_steams)} steams encontrados.")

    # Crear un DataFrame con los steams de cada documento en columnas separadas
    results_df = pd.DataFrame(dict([(k, pd.Series(v)) for k, v in results.items()]))

    # Lectura de Excel de Steams (encontrados en la columna de términos)
    df_steams = pd.read_excel(steams_dir)
    
    # Verifica si la columna 'Steams' existe en el DataFrame df_steams
    if 'Steams' in df_steams.columns:
        # Almacena únicamente la columna 'Steams'
        termino = df_steams[['Steams']]  # Asegúrate de que siga siendo un DataFrame
        
        # Inicializar la matriz booleana
        boolean_matrix = []

        # Para cada término único, verifica su presencia en cada documento
        for term in termino['Steams']:
            row = {'Termino': term}
            for col in results_df.columns:
                row[col] = 1 if term in results_df[col].values else 0
            boolean_matrix.append(row)

        # Convertir la matriz booleana en un DataFrame
        boolean_df = pd.DataFrame(boolean_matrix)

        # Concatenar los DataFrames
        final_df = pd.concat([termino, boolean_df.drop(columns='Termino', errors='ignore')], axis=1)
    else:
        print("La columna 'Steams' no existe en el archivo de steams procesados.")
        return

    # Guardar el DataFrame en un archivo Excel
    output_file_path = os.path.join(output_dir, '6MatrizBinariaDeSteams.xlsx')
    final_df.to_excel(output_file_path, index=False)  # Guardamos en Excel sin encabezados

    print(f"Resultados guardados en {output_file_path}.")

    # Retornar la cantidad total de documentos procesados
    return len(results)

# Ejemplo de uso
corpus_dir = "Corpus"  # Carpeta con los documentos
steams_dir = "PreprocesamientoSteps/4DiccDeSteams.xlsx"  # Archivo Excel con los steams
output_dir = "PreprocesamientoSteps"  # Carpeta para guardar el archivo Excel de salida

# Crear la carpeta de salida si no existe
if not os.path.exists(output_dir):
    os.makedirs(output_dir)

# Procesar los documentos y mostrar el conteo de steams por documento
total_documents = process_and_tokenize_documents(corpus_dir, steams_dir, output_dir)
print(f"Total de documentos procesados: {total_documents}")
