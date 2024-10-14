

import os
import pandas as pd
import unicodedata
from collections import deque
from nltk.tokenize import word_tokenize
from nltk.stem.snowball import SnowballStemmer
import tkinter as tk
from tkinter import ttk, messagebox
from arbol import BTree
import random
from nltk.corpus import stopwords
from nltk.stem import PorterStemmer
import PyPDF2
import docx
from pptx import Presentation

###############  PREPROCESADO DE DOCUMENTOS  ###########################

def read_powerpoint(file_path):
    prs = Presentation(file_path)
    content = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content += shape.text + "\n"
    return content

def read_excel(file_path):
    content = ""
    df = pd.read_excel(file_path)
    for column in df.columns:
        content += " ".join(df[column].astype(str)) + " "
    return content

def read_word(file_path):
    doc = docx.Document(file_path)
    content = "\n".join([para.text for para in doc.paragraphs])
    return content

def read_pdf(file_path):
    content = ""
    with open(file_path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        for page in reader.pages:
            content += page.extract_text()
    return content

def tokenize_and_process_documents(corpus_dir):
    initial_dictionary = {}

    stemmer = PorterStemmer()
    stop_words = set(word.upper() for word in stopwords.words('spanish'))  # Stopwords en mayúsculas

    output_folder = create_folder()  # Crear la carpeta para los archivos de salida

    for filename in os.listdir(corpus_dir):
        file_path = os.path.join(corpus_dir, filename)
        content = ""

        if filename.endswith(".txt"):
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
        elif filename.endswith(".pdf"):
            content = read_pdf(file_path)
        elif filename.endswith(".docx"):
            content = read_word(file_path)
        elif filename.endswith(".xlsx"):
            content = read_excel(file_path)
        elif filename.endswith(".pptx"):
            content = read_powerpoint(file_path)
        else:
            print(f"Ignorando archivo no soportado: {filename}")
            continue
    # Inicio del primer proceso
        tokens = word_tokenize(content)
        tokens = [word for word in tokens if word.isalnum() and not word.isnumeric()]
        tokens = [unicodedata.normalize('NFKD', word).encode('ascii', 'ignore').decode('utf-8') for word in tokens]

        for token in tokens:
            if token not in initial_dictionary:
                initial_dictionary[token] = 1
            else:
                initial_dictionary[token] += 1

    # Guardar el diccionario inicial en "1Diccionario.xlsx"
    df_initial = pd.DataFrame(list(initial_dictionary.items()), columns=['Termino', 'Frecuencia'])
    df_initial.to_excel(os.path.join(output_folder, 'PreprocesamientoSteps/1Diccionario.xlsx'), index=False)
    print('Diccionario inicial guardado en "1Diccionario.xlsx".')


    # inicio del segundo proceso
    upper_dictionary = {}

    # Proceso de cambio a mayúsculas tomando como base el diccionario inicial
    for token in initial_dictionary:
     token_upper = token.upper()
     if token_upper not in upper_dictionary:
         upper_dictionary[token_upper] = initial_dictionary[token]
    else:
         upper_dictionary[token_upper] += initial_dictionary[token]  # Si ya existe, sumar la frecuencia

    # Guardar el diccionario con mayúsculas en "2DiccMayus.xlsx"
    df_upper = pd.DataFrame(list(upper_dictionary.items()), columns=['Termino', 'Frecuencia'])
    df_upper.to_excel(os.path.join(output_folder, 'PreprocesamientoSteps/2DiccMayus.xlsx'), index=False)
    print('Diccionario con mayúsculas guardado en "2DiccMayus.xlsx".')

    #Proceso numero 3
    stop_words = set(stopwords.words('spanish'))
    stop_words_upper = {word.upper() for word in stop_words}
    filtered_dictionary = {word: count for word, count in upper_dictionary.items() if word not in stop_words_upper}

    # Guardar el diccionario sin stop words en "3DiccMSinStopWords.xlsx"
    df_filtered = pd.DataFrame(list(filtered_dictionary.items()), columns=['Termino', 'Frecuencia'])
    df_filtered.to_excel(os.path.join(output_folder, 'PreprocesamientoSteps/3DiccMSinStopWords.xlsx'), index=False)
    print('Diccionario sin stop words guardado en "3DiccMSinStopWords.xlsx".')
    
    #Proceso numnero 4
    stemmer = SnowballStemmer("spanish")
    stemmed_dictionary = {}
    
    for word, count in filtered_dictionary.items():
        stemmed_word = stemmer.stem(word).upper()
        if stemmed_word not in stemmed_dictionary:
            stemmed_dictionary[stemmed_word] = count
        else:
            stemmed_dictionary[stemmed_word] += count

    # Guardar el diccionario con stemming en "4DiccDeSteams.xlsx"
    df_stemmed = pd.DataFrame(list(stemmed_dictionary.items()), columns=['Termino', 'Frecuencia'])
    df_stemmed.to_excel(os.path.join(output_folder, 'PreprocesamientoSteps/4DiccDeSteams.xlsx'), index=False)
    print('Diccionario con stemming guardado en "4DiccDeSteams.xlsx".')

################  PREPROCESADO DE DOCUMENTOS  ##########################



#########################################
def opendocument(link):
    try:
        if os.name == 'nt':  # Para Windows
            os.startfile(link)
    except Exception as e:
        print(f"Error al intentar abrir el documento: {e}")

def preprocess_documents():
    corpus_dir = 'Corpus'
    if os.path.exists(corpus_dir):
        print("Preprocesando documentos...")
        tokenize_and_process_documents(corpus_dir)
    else:
        print("El directorio 'Corpus' no existe.")

def read_steams_from_excel(file_path):
    df = pd.read_excel(file_path)
    steams = df['Termino'].astype(str).tolist()  # Aseguramos que todo sea cadena
    return steams

def index_steams():
    # Leer los steams del archivo generado en el proceso anterior (4DiccDeSteams.xlsx)
    file_path = "PreprocesamientoSteps/4DiccDeSteams.xlsx"
    if not os.path.exists(file_path):
        print(f"El archivo {file_path} no existe.")
        return

    steams = read_steams_from_excel(file_path)

    # Revolver los steams (mezclar aleatoriamente)
    random.shuffle(steams)

    # Crear una tabla hash donde cada steam tendrá un valor hash
    steam_hash_table = {steam: hash(steam) for steam in steams}

    # Guardar la tabla en un archivo Excel
    output_folder = create_folder()
    df = pd.DataFrame(list(steam_hash_table.items()), columns=['Steam', 'Hash'])
    output_file = os.path.join(output_folder, 'PreprocesamientoSteps/5ListDiccIndex.xlsx')
    df.to_excel(output_file, index=False)
    print(f"Tabla hash guardada en {output_file}.")
    return steam_hash_table  # Devolvemos la tabla hash para su uso posterior

def create_folder():
    folder_name = "PreprocesamientoSteps"
    if not os.path.exists(folder_name):
        os.makedirs(folder_name)
        print(f'Carpeta "{folder_name}" creada con éxito.')
    else:
        print(f'La carpeta "{folder_name}" ya existe.')
    return folder_name
#########################################


# Inicializa el stemmer para español
stemmer = SnowballStemmer("spanish")

# Función para convertir la consulta a notación postfija
def infix_to_postfix(consulta):
    precedencia = {'not': 3, 'and': 2, 'or': 1}
    output = []
    operadores = []
    
    for token in consulta:
        if token not in precedencia and token != '(' and token != ')':
            output.append(token)
        elif token == '(':
            operadores.append(token)
        elif token == ')':
            while operadores and operadores[-1] != '(':
                output.append(operadores.pop())
            operadores.pop()  # Quitar '(' de la pila
        else:
            while (operadores and precedencia.get(token, 0) <= precedencia.get(operadores[-1], 0)):
                output.append(operadores.pop())
            operadores.append(token)

    while operadores:
        output.append(operadores.pop())
    
    return output

# Función para tokenizar y obtener la raíz de las palabras del contenido
def tokenize_and_stem(content):
    tokens = word_tokenize(content)
    stemmed_tokens = {stemmer.stem(token).upper() for token in tokens if token.isalpha() and len(token) > 1}
    return sorted(stemmed_tokens)
# Cargar el Excel
archivo_excel = 'PreprocesamientoSteps/6MatrizBinariaDeSteams.xlsx'
df = pd.read_excel(archivo_excel, index_col=0)

hash_table = {col: list(df.loc[col].values) for col in df.index}

# Función para evaluar la consulta en notación postfija
def evaluate_postfix(postfix, hash_table):
    stack = deque()

    for token in postfix:
        # Aseguramos que los operadores no sean considerados palabras
        if token not in {'and', 'or', 'not'}:
            # Comparamos con las palabras ya procesadas (en mayúsculas)
            if token in hash_table:
                term_vector = hash_table[token]
                stack.append(deque(term_vector))
            else:
                print(f"Advertencia: El término '{token}' no se encuentra en la tabla hash.")
                stack.append(deque([0] * len(next(iter(hash_table.values())))))
        else:
            if token == 'not':
                a = stack.pop()
                print(f"Evaluando 'not': {list(a)}")  # Mostrar vector antes de aplicar 'not'
                stack.append(deque([1 - x for x in a]))
            else:
                b = stack.pop()
                a = stack.pop()
                if token == 'and':
                    print(f"Evaluando 'and': {list(a)} AND {list(b)}")  # Mostrar vectores antes de aplicar 'and'
                    stack.append(deque([x & y for x, y in zip(a, b)]))
                elif token == 'or':
                    print(f"Evaluando 'or': {list(a)} OR {list(b)}")  # Mostrar vectores antes de aplicar 'or'
                    stack.append(deque([x | y for x, y in zip(a, b)]))
    
    return list(stack.pop())

# Cargar el Excel
archivo_excel = 'PreprocesamientoSteps/7MatrizFrecuenciasDeSteams.xlsx'
df = pd.read_excel(archivo_excel, index_col=0)

hash_tableS = {col: list(df.loc[col].values) for col in df.index}

# Función para evaluar la consulta en notación postfija
def evaluate_postfi_ex(postfix, hash_tableS):
    stack = deque()  # Pila para evaluar la expresión postfija

    for token in postfix:
        if token not in {'and', 'or', 'not', '(', ')'}:  # Ignorar paréntesis y operadores
            # Agregar el término de la tabla hash a la pila
            if token in hash_tableS:
                stack.append(hash_tableS[token])  # Colocar tf(x) en la pila
            else:
                print(f"Término no encontrado en la tabla hash: {token}")
                # Si el término no existe, usa una lista de ceros del tamaño adecuado
                stack.append(deque([0] * len(next(iter(hash_tableS.values())))))
        else:
            # Condición: NOT x -> 0 si tf(x) > 0, o 1 si tf(x) == 0
            if token == 'not':
                a = stack.pop()
                # Aplicar la regla para NOT x
                stack.append(deque([0 if x > 0 else 1 for x in a]))  # tf(x) > 0 -> 0; tf(x) == 0 -> 1
            else:
                b = stack.pop()  # Segundo operando (para AND/OR)
                a = stack.pop()  # Primer operando (para AND/OR)

                if token == 'and':
                    # Condición: x AND y -> tf(x) * tf(y)
                    stack.append(deque([x * y for x, y in zip(a, b)]))  # Multiplicar tf(x) y tf(y)
                
                elif token == 'or':
                    # Condición: x OR y -> tf(x) + tf(y)
                    stack.append(deque([x + y for x, y in zip(a, b)]))  # Sumar tf(x) y tf(y)

    # Retornar el resultado final de la pila
    return list(stack.pop())

class App:
    def __init__(self, root):
        self.root = root
        self.root.title("Cronus con Árbol B")
        self.root.configure(bg='black')
        self.root.attributes('-fullscreen', True)
        
        self.b_tree = BTree(t=4)  # Crear un árbol B de grado 4
        self.hash_table = {}  # Tabla hash vacía
        self.create_widgets()

        # Cargar el Excel de matriz binaria
        archivo_excel = 'PreprocesamientoSteps/6MatrizBinariaDeSteams.xlsx'
        self.df = pd.read_excel(archivo_excel, index_col=0)  
        self.hash_table = {col: list(self.df.loc[col].values) for col in self.df.index}
        print("Matriz binaria de documentos (cargada desde Excel):")
        print(self.df)

    def create_widgets(self):
        # Estilo
        style = ttk.Style()
        style.configure("TButton", font=("arial", 16))
        style.configure("TEntry", font=("arial", 16))

        #PRIMER cuadro de texto para busqueda del nombre del archivo
        search_frame = tk.Frame(self.root, bg='black')
        search_frame.pack(pady=20)

        self.search_entry = ttk.Entry(search_frame, width=30, style="TEntry")
        self.search_entry.pack(side=tk.LEFT, padx=5)

        search_button = ttk.Button(search_frame, text="🔍", command=self.search_word, style="TButton")
        search_button.pack(side=tk.LEFT)

        # SEGUNDO cuadro de texto para consulta booleana
        consulta_frame = tk.Frame(self.root, bg='black')
        consulta_frame.pack(pady=20)

        self.consulta_entry = ttk.Entry(consulta_frame, width=50, style="TEntry")
        self.consulta_entry.pack(side=tk.LEFT, padx=5)

        # Botón para procesar la consulta booleana
        consulta_button = ttk.Button(consulta_frame, text="Procesar consulta booleana", command=self.process_query, style="TButton")
        consulta_button.pack(side=tk.LEFT)

        # TERCER cuadro de texto para consulta booleana
        consulta_extentds = tk.Frame(self.root, bg='black')
        consulta_extentds.pack(pady=20)

        self.consulta_str_extends = ttk.Entry(consulta_extentds, width=50, style="TEntry")
        self.consulta_str_extends.pack(side=tk.LEFT, padx=5)

        # Botón para procesar la consulta booleana
        button_extentds = ttk.Button(consulta_extentds, text="Procesar consulta booleana extends", command=self.process_query_extends, style="TButton")
        button_extentds.pack(side=tk.LEFT)
        
        # Botón para preprocesar documentos
        preprocess_button = ttk.Button(self.root, text="Preprocesado", command=preprocess_documents, style="TButton")
        preprocess_button.pack(pady=10)

        # Botón para indexar los steams
        index_button = ttk.Button(self.root, text="Indexación", command=self.index_steams, style="TButton")
        index_button.pack(pady=10)

        # Botón para cargar steams en el Árbol B
        load_button = ttk.Button(self.root, text="Cargar STEAMS", command=self.load_steams, style="TButton")
        load_button.pack(pady=10)

        # Etiqueta para mostrar el resultado de la búsqueda
        self.result_label = tk.Label(self.root, text="", fg="white", bg="black", font=("arial", 16))
        self.result_label.pack(pady=10)

        # Canvas para mostrar documentos del corpus
        canvas = tk.Canvas(self.root, bg='gray')
        canvas.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

        scrollbar = ttk.Scrollbar(self.root, orient="vertical", command=canvas.yview)
        scrollbar.pack(side=tk.RIGHT, fill='y')

        canvas.configure(yscrollcommand=scrollbar.set)
        canvas.bind('<Configure>', lambda e: canvas.configure(scrollregion=canvas.bbox("all")))

        frame = tk.Frame(canvas, bg='gray')
        canvas.create_window((0, 0), window=frame, anchor='nw')

        corpus_dir = 'Corpus'
        if not os.path.exists(corpus_dir):
            os.makedirs(corpus_dir)

        documents = [f for f in os.listdir(corpus_dir) if os.path.isfile(os.path.join(corpus_dir, f))]

        for doc in documents:
            doc_path = os.path.join(corpus_dir, doc)
            link = tk.Label(frame, text=doc, fg="white", cursor="hand2", bg='gray', font=("arial", 16))
            link.pack(anchor='w', pady=2)
            link.bind("<Button-1>", lambda e, url=doc_path: opendocument(url))

        self.root.bind("<Escape>", lambda e: self.root.attributes("-fullscreen", False))

    def load_steams(self):
        # Leer el archivo de Excel
        file_path = "PreprocesamientoSteps/4DiccDeSteams.xlsx"
        if not os.path.exists(file_path):
            messagebox.showerror("Error", f"El archivo {file_path} no existe.")
            return
        
        steams = read_steams_from_excel(file_path)
        for steam in steams:
            try:
                self.b_tree.insert(steam)
            except Exception as e:
                messagebox.showerror("Error de inserción", f"Error al insertar '{steam}': {str(e)}")

        self.b_tree.display()  # Mostrar el árbol en consola
        messagebox.showinfo("Éxito", "STEAMS cargados en el árbol.")

    def search_word(self):
        word = self.search_entry.get()
        if word:
            found = self.b_tree.search(word)
            if found:
                self.result_label.config(text=f"La palabra '{word}' fue encontrada en el árbol.")
            else:
                self.result_label.config(text=f"La palabra '{word}' NO fue encontrada en el árbol.")
        else:
            messagebox.showwarning("Entrada vacía", "Por favor, ingresa una palabra para buscar.")

    def index_steams(self):
        self.hash_table = index_steams()
        if self.hash_table:
            messagebox.showinfo("Éxito", "Indexación completada.")
            
    def process_query_extends(self):
        #Obtener la consultad booleanda extendida
        consulta_estends = self.consulta_str_extends.get()
        if consulta_estends:
            print(consulta_estends)
            
            try:
                print(f"Palabras tokenizadas y en su raíz: {tokens}")

                # Convertir la consulta a formato infijo
                consulta = consulta_estends.replace('(', ' ( ').replace(')', ' ) ').split()
                
                # Convertir la consulta infija a postfija
                consulta_postfija = infix_to_postfix(consulta)
                
                # Convertir los tokens a su raíz para la notación postfija
                consulta_postfija = [stemmer.stem(token).upper() if token not in {'(', ')', 'and', 'or', 'not'} else token for token in consulta_postfija]
                print(f"Consulta en notación postfija: {consulta_postfija}")

                # Evaluar la consulta en notación postfija
                resultado = evaluate_postfi_ex(consulta_postfija, self.hash_table)
                
                # Crear el DataFrame sin el primer elemento del resultado (resultado[1:]) y sin el índice
                resultado_matriz = pd.DataFrame([resultado[1:]], columns=self.df.columns[1:])  # Ignorar el primer elemento y ajustar las columnas

                # Mostrar el resultado en la interfaz
                self.result_label.config(text=f"Resultado de la consulta:\n{resultado_matriz}")

                # Guardar el resultado en Excel si es necesario
                ruta_resultado = 'challengeExample/resultado_consulta_ext.xlsx'
                resultado_matriz.to_excel(ruta_resultado, index=False)
                print(f"Resultado guardado en el archivo Excel: {ruta_resultado}")

            except ValueError as e:
                messagebox.showerror("Error en la evaluación de la consulta", str(e))
        else:
            messagebox.showwarning("Consulta vacía", "Por favor, ingresa una consulta booleana.") 
            
        
    def process_query(self):
        # Obtener la consulta booleana ingresada
        consulta_str = self.consulta_entry.get()
        if consulta_str:
            try:
                # Tokenizar y obtener los steams de la consulta ingresada
                tokens = tokenize_and_stem(consulta_str)
                print(f"Palabras tokenizadas y en su raíz: {tokens}")

                # Convertir la consulta a formato infijo
                consulta = consulta_str.replace('(', ' ( ').replace(')', ' ) ').split()
                
                # Convertir la consulta infija a postfija
                consulta_postfija = infix_to_postfix(consulta)
                
                # Convertir los tokens a su raíz para la notación postfija
                consulta_postfija = [stemmer.stem(token).upper() if token not in {'(', ')', 'and', 'or', 'not'} else token for token in consulta_postfija]
                print(f"Consulta en notación postfija: {consulta_postfija}")

                # Evaluar la consulta en notación postfija
                resultado = evaluate_postfix(consulta_postfija, self.hash_table)
                
                # Crear el DataFrame sin el primer elemento del resultado (resultado[1:]) y sin el índice
                resultado_matriz = pd.DataFrame([resultado[1:]], columns=self.df.columns[1:])  # Ignorar el primer elemento y ajustar las columnas

                # Mostrar el resultado en la interfaz
                self.result_label.config(text=f"Resultado de la consulta:\n{resultado_matriz}")

                # Guardar el resultado en Excel si es necesario
                ruta_resultado = 'challengeExample/resultado_consulta_bool.xlsx'
                resultado_matriz.to_excel(ruta_resultado, index=False)
                print(f"Resultado guardado en el archivo Excel: {ruta_resultado}")

            except ValueError as e:
                messagebox.showerror("Error en la evaluación de la consulta", str(e))
        else:
            messagebox.showwarning("Consulta vacía", "Por favor, ingresa una consulta booleana.") 


if __name__ == "__main__":
    root = tk.Tk()
    app = App(root)
    root.mainloop()